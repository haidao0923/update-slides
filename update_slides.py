from pptx import Presentation
from pptx.util import Inches
import win32com.client
import os
DATE_INDEX = 2
TIME_INDEX = 3
LOCATION_INDEX = 4
text_array = []

from PIL import Image

def get_all_text(powerpoint_name, presentation):
    global text_array
    text_array = []
    slide = presentation.slides[0]
    for element in slide.shapes:
        if element.has_text_frame:
            text_array.append(element);


def modify_text(index, new_text):
    text_frame = text_array[index].text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.runs[0].text = new_text

def save(powerpoint_name, presentation, image_folder, image_name):
    presentation.save(powerpoint_name)
    powerpoint_path = os.getcwd() + f'/{powerpoint_name}'
    Application = win32com.client.Dispatch('Powerpoint.Application')
    powerpoint = Application.Presentations.Open(powerpoint_path)
    image_path = os.getcwd() + f'/{image_folder}/{image_name}'
    Application.ActivePresentation.Slides[0].Export(image_path, 'PNG')
    powerpoint.Close()
    Application.Quit()
    os.system('taskkill /F /IM POWERPNT.EXE')

def execute(powerpoint_name, image_folder, image_name, date_text, time_text, location_text):
    presentation = Presentation(powerpoint_name)
    get_all_text(powerpoint_name, presentation)
    for i in range(len(text_array)):
        print(str(i) + text_array[i].text)
    modify_text(DATE_INDEX, date_text)
    modify_text(TIME_INDEX, time_text)
    modify_text(LOCATION_INDEX, location_text)
    save(powerpoint_name, presentation, image_folder, image_name);

wednesday_date = ["09_27_23"]
sunday_date = ["10_01_23"]
wednesday_location_list = ['Student Center Peachtree']
sunday_location_list = ['Student Center Northside']

for i in range(len(wednesday_date)):
    image_date = sunday_date[i]
    text_date = image_date[0:2] + "/" + image_date[3:5] + "/" + image_date[6:8]
    location = sunday_location_list[i]
    wednesday_location = wednesday_location_list[i]
    jennaral_meeting_date = wednesday_date[i]
    jennaral_text_date = jennaral_meeting_date[0:2] + "/" + jennaral_meeting_date[3:5] + "/" + jennaral_meeting_date[6:8]
    #online_image_date = '01_04_23'
    #online_text_date = '01/04/23'

    execute('GGTemplate.pptx', 'Weekly_Images', f'Engage_{image_date}.png',
            f'Sunday {text_date}', '2pm - 5pm', location)
    execute('GGTemplateSquare.pptx', 'Weekly_Images', f'Instagram_{image_date}.png',
            f'Sunday {text_date}', '2pm - 5pm', location)
    execute('JennaralMeetingTemplate.pptx', 'Weekly_Images', f'Engage_Jennaral_{jennaral_meeting_date}.png',
            f'Wednesday {jennaral_text_date}', '7:30pm - 10:30pm', wednesday_location)
    execute('JennaralMeetingTemplateSquare.pptx', 'Weekly_Images', f'Instagram_Jennaral_{jennaral_meeting_date}.png',
            f'Wednesday {jennaral_text_date}', '7:30pm - 10:30pm', wednesday_location)
    #execute('GGOnlineTemplate.pptx', 'Weekly_Images', f'Engage_Online_{online_image_date}.png',
    #        f'Wednesday {online_text_date}', '9pm', 'Discord Voice Chat - Online')
    #execute('GGOnlineTemplateSquare.pptx', 'Weekly_Images', f'Instagram_Online_{online_image_date}.png',
    #        f'Wednesday {online_text_date}', '9pm', 'Discord Voice Chat - Online')
    first_image = Image.open(os.getcwd() + '/Weekly_Images/' + f'Instagram_Jennaral_{jennaral_meeting_date}.png')
    display = Image.new('RGB', (first_image.width * 2, first_image.height))
    display.paste(first_image, (0, 0))
    second_image = Image.open(os.getcwd() + '/Weekly_Images/' + f'Instagram_{image_date}.png')
    display.paste(second_image, (first_image.width, 0))
    display.save(os.getcwd() + '/Weekly_Images/' + f'Combined_Instagram_{image_date}.png')
    combined_image = Image.open(os.getcwd() + '/Weekly_Images/' + f'Combined_Instagram_{image_date}.png')
    bordered_display = Image.new('RGB', (combined_image.width, combined_image.width))
    bordered_display.paste(combined_image, (0, int((combined_image.width - combined_image.height) / 2)))
    bordered_display.save(os.getcwd() + '/Weekly_Images/' + f'Bordered_Instagram_{image_date}.png')
